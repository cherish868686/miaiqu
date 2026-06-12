"""
PPT文档解析器 (.pptx)
提取每页幻灯片的文本和表格内容
"""
from pptx import Presentation


def parse_pptx(filepath):
    """解析PPT文档，返回文本和表格内容"""
    prs = Presentation(filepath)
    texts = []
    tables = []

    for slide_num, slide in enumerate(prs.slides, 1):
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        texts.append({
                            'index': len(texts),
                            'content': text,
                            'slide': slide_num
                        })

            if shape.has_table:
                table = shape.table
                table_data = {
                    'index': len(tables),
                    'slide': slide_num,
                    'headers': [],
                    'rows': []
                }
                for r_idx, row in enumerate(table.rows):
                    row_data = [cell.text.strip() for cell in row.cells]
                    if r_idx == 0:
                        table_data['headers'] = row_data
                    else:
                        table_data['rows'].append({
                            'row_index': r_idx,
                            'cells': row_data
                        })
                tables.append(table_data)

    return {'texts': texts, 'tables': tables}
